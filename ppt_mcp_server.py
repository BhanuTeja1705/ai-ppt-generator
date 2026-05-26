from mcp.server.fastmcp import FastMCP
from pptx import Presentation

from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

# =========================================
# MCP SERVER
# =========================================

mcp = FastMCP("Modern PPT MCP Server")

prs = None

# =========================================
# THEME
# =========================================

THEME = {
    "bg": RGBColor(15, 23, 42),
    "card": RGBColor(30, 41, 59),
    "title": RGBColor(255, 255, 255),
    "text": RGBColor(220, 220, 220),
    "accent": RGBColor(99, 102, 241)
}

ACCENTS = [
    RGBColor(99, 102, 241),
    RGBColor(16, 185, 129),
    RGBColor(236, 72, 153),
    RGBColor(245, 158, 11),
]

# =========================================
# CREATE PRESENTATION
# =========================================

@mcp.tool()
def create_presentation():

    global prs

    prs = Presentation()

    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    return "Presentation Created"

# =========================================
# TITLE SLIDE
# =========================================

@mcp.tool()
def add_title_slide(title: str, subtitle: str):

    global prs

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["bg"]

    # Left Accent Bar
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(0.25),
        Inches(7.5)
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME["accent"]

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(10),
        Inches(1.5)
    )

    tf = title_box.text_frame

    p = tf.paragraphs[0]

    p.text = title

    p.font.size = Pt(42)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = THEME["title"]

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(1),
        Inches(3.3),
        Inches(8),
        Inches(1)
    )

    sub_tf = sub_box.text_frame

    sub_p = sub_tf.paragraphs[0]

    sub_p.text = subtitle

    sub_p.font.size = Pt(20)
    sub_p.font.name = "Calibri"
    sub_p.font.color.rgb = THEME["text"]

    return "Title Slide Added"

# =========================================
# CONTENT SLIDE
# =========================================

@mcp.tool()
def add_content_slide(title: str, points: list):

    global prs

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["bg"]

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.4),
        Inches(12),
        Inches(0.8)
    )

    tf = title_box.text_frame

    p = tf.paragraphs[0]

    p.text = title

    p.font.size = Pt(34)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = THEME["title"]

    # Accent Line
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.6),
        Inches(1.12),
        Inches(2.5),
        Inches(0.07)
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME["accent"]

    # Content Box
    content_box = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(1.45),
        Inches(12),
        Inches(5.5)
    )

    content_tf = content_box.text_frame

    content_tf.clear()

    for point in points:

        para = content_tf.add_paragraph()

        para.text = "▸ " + point

        para.font.size = Pt(24)
        para.font.name = "Calibri"
        para.font.color.rgb = THEME["text"]

        para.space_after = Pt(24)

        para.level = 0

    return "Content Slide Added"

# =========================================
# FULL WIDTH CARDS SLIDE
# =========================================

@mcp.tool()
def add_cards_slide(title: str, points: list):

    global prs

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["bg"]

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.4),
        Inches(12),
        Inches(0.8)
    )

    tf = title_box.text_frame

    p = tf.paragraphs[0]

    p.text = title

    p.font.size = Pt(34)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = THEME["title"]

    # Accent Line
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.6),
        Inches(1.12),
        Inches(2.5),
        Inches(0.07)
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME["accent"]

    # Responsive Layout
    card_count = min(len(points), 6)

    if card_count <= 3:

        positions = [
            (0.45, 2.1),
            (4.45, 2.1),
            (8.15, 2.1),
        ]

        width = 3.75
        height = 2.7

    else:

        positions = [
            (0.45, 1.7),
            (4.3, 1.7),
            (8.15, 1.7),

            (0.45, 4.25),
            (4.3, 4.25),
            (8.15, 4.25),
        ]

        width = 3.75
        height = 2.05

    # Create Cards
    for i, point in enumerate(points[:6]):

        x, y = positions[i]

        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(width),
            Inches(height)
        )

        card.fill.solid()
        card.fill.fore_color.rgb = THEME["card"]

        card.line.color.rgb = ACCENTS[i % len(ACCENTS)]
        card.line.width = Pt(2)

        tf = card.text_frame

        tf.clear()

        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]

        p.text = point

        p.font.size = Pt(18)
        p.font.bold = True
        p.font.name = "Calibri"
        p.font.color.rgb = THEME["title"]

        p.alignment = PP_ALIGN.CENTER

    return "Cards Slide Added"

# =========================================
# TIMELINE SLIDE
# =========================================

@mcp.tool()
def add_timeline_slide(title: str, steps: list):

    global prs

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["bg"]

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.4),
        Inches(12),
        Inches(0.8)
    )

    tf = title_box.text_frame

    p = tf.paragraphs[0]

    p.text = title

    p.font.size = Pt(34)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = THEME["title"]

    # Accent Line
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.6),
        Inches(1.12),
        Inches(2.5),
        Inches(0.07)
    )

    accent.fill.solid()
    accent.fill.fore_color.rgb = THEME["accent"]

    # Timeline Line
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(1),
        Inches(3.7),
        Inches(10.5),
        Inches(0.08)
    )

    line.fill.solid()
    line.fill.fore_color.rgb = THEME["accent"]

    x = 1

    for i, step in enumerate(steps[:5]):

        circle = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x),
            Inches(3.35),
            Inches(0.7),
            Inches(0.7)
        )

        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENTS[i % len(ACCENTS)]

        # Step Number
        num_box = slide.shapes.add_textbox(
            Inches(x + 0.22),
            Inches(3.43),
            Inches(0.2),
            Inches(0.2)
        )

        num_tf = num_box.text_frame

        num_p = num_tf.paragraphs[0]

        num_p.text = str(i + 1)

        num_p.font.size = Pt(10)
        num_p.font.bold = True
        num_p.font.color.rgb = RGBColor(255,255,255)

        # Step Text
        box = slide.shapes.add_textbox(
            Inches(x - 0.35),
            Inches(2),
            Inches(1.9),
            Inches(1)
        )

        tf = box.text_frame

        p = tf.paragraphs[0]

        p.text = step

        p.font.size = Pt(14)
        p.font.name = "Calibri"
        p.font.color.rgb = THEME["text"]

        p.alignment = PP_ALIGN.CENTER

        x += 2.25

    return "Timeline Slide Added"

# =========================================
# SAVE PPT
# =========================================

@mcp.tool()
def save_presentation(filename: str):

    global prs

    prs.save(filename)

    return f"Saved as {filename}"

# =========================================
# RUN SERVER
# =========================================

if __name__ == "__main__":
    mcp.run()
